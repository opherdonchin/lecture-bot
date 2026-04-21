from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest
from pptx import Presentation

ROOT_DIR = Path(__file__).resolve().parents[1]
if str(ROOT_DIR) not in sys.path:
    sys.path.insert(0, str(ROOT_DIR))

from scripts import build_lecture_package


def write_lecture_config(lecture_dir: Path) -> None:
    config = {
        "lecture_id": lecture_dir.name,
        "title": "Fixture Lecture",
        "course": "Test Course",
        "active": True,
        "files": {
            "slides": {"source": "slides_source.pptx", "target": "slides.md"},
            "handout": {"source": "handout_source.qmd", "target": "handout.md"},
            "notebook": {"source": "notebook_source.ipynb", "target": "notebook.md"},
            "rubric": {"source": "rubric_source.md", "target": "rubric.md"},
        },
    }
    lecture_dir.joinpath("lecture_config.json").write_text(
        json.dumps(config, indent=2) + "\n",
        encoding="utf-8",
    )


def write_generated_pipeline_config(lecture_dir: Path) -> None:
    config = {
        "lecture_id": lecture_dir.name,
        "title": "Fixture Lecture",
        "course": "Test Course",
        "active": True,
        "files": {
            "slides": {"source": "slides_source.pptx", "target": "slides.md"},
            "handout": {"source": "handout_source.qmd", "target": "handout.md"},
            "notebook": {"source": "notebook_source.ipynb", "target": "notebook.md"},
            "transcript": {"source": "lecture_transcript.vtt", "target": "transcript.md"},
            "minutes": {"target": "minutes.json"},
            "rubric": {"target": "rubric.md"},
        },
    }
    lecture_dir.joinpath("lecture_config.json").write_text(
        json.dumps(config, indent=2) + "\n",
        encoding="utf-8",
    )


def write_placeholder_sources(lecture_dir: Path) -> None:
    lecture_dir.joinpath("slides_source.pptx").write_text("placeholder", encoding="utf-8")
    lecture_dir.joinpath("handout_source.qmd").write_text("# Handout\n", encoding="utf-8")
    lecture_dir.joinpath("notebook_source.ipynb").write_text("{}", encoding="utf-8")
    lecture_dir.joinpath("rubric_source.md").write_text("Rubric\n", encoding="utf-8")


def write_transcript_source(lecture_dir: Path) -> None:
    lecture_dir.joinpath("lecture_transcript.vtt").write_text(
        "WEBVTT\n\n"
        "00:00:01.000 --> 00:00:04.000\n"
        "Welcome to Bayesian models.\n\n"
        "00:00:04.000 --> 00:00:07.000\n"
        "<v Professor>Today we connect prior and likelihood.\n",
        encoding="utf-8",
    )


def write_real_sources(lecture_dir: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Probability Basics"
    text_frame = slide.placeholders[1].text_frame
    text_frame.text = "Outcome space"
    bullet = text_frame.add_paragraph()
    bullet.text = "Events"
    bullet.level = 1
    presentation.save(lecture_dir / "slides_source.pptx")

    lecture_dir.joinpath("handout_source.qmd").write_text(
        "---\n"
        "title: Fixture\n"
        "---\n\n"
        "# Handout\n\n"
        "Quarto content.\n",
        encoding="utf-8",
    )

    notebook = {
        "cells": [
            {
                "cell_type": "markdown",
                "metadata": {},
                "source": ["# Notebook heading\n", "\n", "Some notes.\n"],
            },
            {
                "cell_type": "code",
                "execution_count": 1,
                "metadata": {},
                "source": ["print('hello world')\n"],
                "outputs": [
                    {
                        "output_type": "stream",
                        "name": "stdout",
                        "text": ["hello world\n"],
                    }
                ],
            },
        ],
        "metadata": {"language_info": {"name": "python"}},
        "nbformat": 4,
        "nbformat_minor": 5,
    }
    lecture_dir.joinpath("notebook_source.ipynb").write_text(
        json.dumps(notebook, indent=2) + "\n",
        encoding="utf-8",
    )

    lecture_dir.joinpath("rubric_source.md").write_text(
        "Rubric content\n",
        encoding="utf-8",
    )


def make_lecture_dir(tmp_path: Path, name: str = "lecture_fixture") -> Path:
    lecture_dir = tmp_path / "lectures" / name
    lecture_dir.mkdir(parents=True)
    write_lecture_config(lecture_dir)
    return lecture_dir


def test_loads_config_and_dispatches_correctly(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    lecture_dir = make_lecture_dir(tmp_path)
    write_placeholder_sources(lecture_dir)

    config = build_lecture_package.load_lecture_config(lecture_dir)
    assert config["lecture_id"] == "lecture_fixture"

    calls: list[tuple[str, str, str]] = []

    def fake_slides(source: Path, target: Path) -> None:
        calls.append(("slides", source.name, target.name))
        target.write_text("slides output\n", encoding="utf-8")

    def fake_handout(source: Path, target: Path) -> None:
        calls.append(("handout", source.name, target.name))
        target.write_text("handout output\n", encoding="utf-8")

    def fake_notebook(source: Path, target: Path) -> None:
        calls.append(("notebook", source.name, target.name))
        target.write_text("notebook output\n", encoding="utf-8")

    monkeypatch.setattr(build_lecture_package.convert_pptx, "convert_pptx_to_md", fake_slides)
    monkeypatch.setattr(build_lecture_package.convert_qmd, "convert_qmd_to_md", fake_handout)
    monkeypatch.setattr(build_lecture_package.convert_ipynb, "convert_ipynb_to_md", fake_notebook)

    build_lecture_package.build_lecture_package(
        "lecture_fixture",
        lectures_root=tmp_path / "lectures",
        force=True,
    )

    assert calls == [
        ("slides", "slides_source.pptx", "slides.md"),
        ("handout", "handout_source.qmd", "handout.md"),
        ("notebook", "notebook_source.ipynb", "notebook.md"),
    ]
    assert lecture_dir.joinpath("rubric.md").read_text(encoding="utf-8") == "Rubric\n"


def test_refuses_to_overwrite_targets_without_force(tmp_path: Path) -> None:
    lecture_dir = make_lecture_dir(tmp_path)
    write_placeholder_sources(lecture_dir)
    lecture_dir.joinpath("slides.md").write_text("existing slides\n", encoding="utf-8")

    with pytest.raises(FileExistsError, match="slides.md"):
        build_lecture_package.build_lecture_package(lecture_dir)

    assert lecture_dir.joinpath("slides.md").read_text(encoding="utf-8") == "existing slides\n"


def test_build_creates_processed_markdown_targets(tmp_path: Path) -> None:
    lecture_dir = make_lecture_dir(tmp_path)
    write_real_sources(lecture_dir)

    jobs = build_lecture_package.build_lecture_package(lecture_dir, force=True)

    assert [job.logical_key for job in jobs] == ["slides", "handout", "notebook", "rubric"]

    slides_text = lecture_dir.joinpath("slides.md").read_text(encoding="utf-8")
    assert "## Slide 1" in slides_text
    assert "### Probability Basics" in slides_text
    assert "Outcome space" in slides_text
    assert "Events" in slides_text

    handout_text = lecture_dir.joinpath("handout.md").read_text(encoding="utf-8")
    assert "title: Fixture" not in handout_text
    assert "# Handout" in handout_text
    assert "Quarto content." in handout_text

    notebook_text = lecture_dir.joinpath("notebook.md").read_text(encoding="utf-8")
    assert "## Markdown Cell 1" in notebook_text
    assert "# Notebook heading" in notebook_text
    assert "## Code Cell 2" in notebook_text
    assert "```python" in notebook_text
    assert "print('hello world')" in notebook_text
    assert "### Output" in notebook_text
    assert "hello world" in notebook_text

    rubric_text = lecture_dir.joinpath("rubric.md").read_text(encoding="utf-8")
    assert rubric_text == "Rubric content\n"


def test_errors_when_source_file_is_missing(tmp_path: Path) -> None:
    lecture_dir = make_lecture_dir(tmp_path)
    write_placeholder_sources(lecture_dir)
    lecture_dir.joinpath("notebook_source.ipynb").unlink()

    with pytest.raises(FileNotFoundError, match="notebook_source.ipynb"):
        build_lecture_package.build_lecture_package(lecture_dir, force=True)


def test_generated_minutes_require_transcript_config(tmp_path: Path) -> None:
    lecture_dir = make_lecture_dir(tmp_path)
    write_placeholder_sources(lecture_dir)

    config = json.loads(lecture_dir.joinpath("lecture_config.json").read_text(encoding="utf-8"))
    config["files"]["minutes"] = {"target": "minutes.json"}
    lecture_dir.joinpath("lecture_config.json").write_text(
        json.dumps(config, indent=2) + "\n",
        encoding="utf-8",
    )

    with pytest.raises(ValueError, match="Generated minutes require a transcript file"):
        build_lecture_package.build_lecture_package(lecture_dir, force=True)


def test_build_generates_minutes_and_rubric_from_transcript_pipeline(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    lecture_dir = tmp_path / "lectures" / "lecture_generated"
    lecture_dir.mkdir(parents=True)
    write_generated_pipeline_config(lecture_dir)
    write_placeholder_sources(lecture_dir)
    write_transcript_source(lecture_dir)

    calls: list[tuple[str, str]] = []

    def fake_slides(source: Path, target: Path) -> None:
        calls.append(("slides", target.name))
        target.write_text("# Slides\n", encoding="utf-8")

    def fake_handout(source: Path, target: Path) -> None:
        calls.append(("handout", target.name))
        target.write_text("# Handout\n", encoding="utf-8")

    def fake_notebook(source: Path, target: Path) -> None:
        calls.append(("notebook", target.name))
        target.write_text("# Notebook\n", encoding="utf-8")

    def fake_transcript(source: Path, target: Path) -> None:
        calls.append(("transcript", target.name))
        target.write_text("# Transcript\n", encoding="utf-8")

    def fake_minutes(*, source_paths: dict[str, Path], target: Path) -> None:
        calls.append(("minutes", target.name))
        assert source_paths["slides"].name == "slides.md"
        assert source_paths["handout"].name == "handout.md"
        assert source_paths["notebook"].name == "notebook.md"
        assert source_paths["transcript"].name == "transcript.md"
        target.write_text('{"lecture_metadata": {"title": "Fixture"}}\n', encoding="utf-8")

    def fake_rubric(*, source_paths: dict[str, Path], target: Path) -> None:
        calls.append(("rubric", target.name))
        assert source_paths["minutes"].name == "minutes.json"
        target.write_text(
            "# Mastery Rubric\n\n"
            "### T1. Generated Topic\n"
            "**Importance:** core\n",
            encoding="utf-8",
        )

    monkeypatch.setattr(build_lecture_package.convert_pptx, "convert_pptx_to_md", fake_slides)
    monkeypatch.setattr(build_lecture_package.convert_qmd, "convert_qmd_to_md", fake_handout)
    monkeypatch.setattr(build_lecture_package.convert_ipynb, "convert_ipynb_to_md", fake_notebook)
    monkeypatch.setattr(build_lecture_package.convert_vtt, "convert_vtt_to_md", fake_transcript)
    monkeypatch.setattr(
        build_lecture_package.generate_artifacts,
        "generate_instructional_minutes",
        fake_minutes,
    )
    monkeypatch.setattr(
        build_lecture_package.generate_artifacts,
        "generate_master_rubric",
        fake_rubric,
    )

    jobs = build_lecture_package.build_lecture_package(
        lecture_dir,
        force=True,
    )

    assert [job.logical_key for job in jobs] == [
        "slides",
        "handout",
        "notebook",
        "transcript",
        "minutes",
        "rubric",
    ]
    assert calls == [
        ("slides", "slides.md"),
        ("handout", "handout.md"),
        ("notebook", "notebook.md"),
        ("transcript", "transcript.md"),
        ("minutes", "minutes.json"),
        ("rubric", "rubric.md"),
    ]
    assert lecture_dir.joinpath("minutes.json").exists()
    assert lecture_dir.joinpath("rubric.md").exists()

    config = json.loads(lecture_dir.joinpath("lecture_config.json").read_text(encoding="utf-8"))
    assert config["topics"] == [
        {
            "topic_id": "T1",
            "label": "Generated Topic",
            "importance": "core",
        }
    ]
