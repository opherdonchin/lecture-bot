from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation


def collapse_text(text: str) -> str:
    return " ".join(part.strip() for part in text.splitlines() if part.strip()).strip()


def text_frame_lines(text_frame) -> list[str]:
    paragraphs = [paragraph for paragraph in text_frame.paragraphs if collapse_text(paragraph.text)]
    if not paragraphs:
        return []

    bullet_mode = len(paragraphs) > 1 or any(paragraph.level > 0 for paragraph in paragraphs)
    lines: list[str] = []

    for paragraph in paragraphs:
        text = collapse_text(paragraph.text)
        if not text:
            continue

        if bullet_mode:
            indent = "  " * paragraph.level
            lines.append(f"{indent}- {text}")
        else:
            lines.append(text)

    return lines


def shape_lines(shape) -> list[str]:
    if getattr(shape, "has_text_frame", False):
        return text_frame_lines(shape.text_frame)

    if getattr(shape, "has_table", False):
        rows: list[str] = []
        for row in shape.table.rows:
            cells = [collapse_text(cell.text) for cell in row.cells]
            cells = [cell for cell in cells if cell]
            if cells:
                rows.append(" | ".join(cells))
        return rows

    child_shapes = getattr(shape, "shapes", None)
    if child_shapes is not None:
        lines: list[str] = []
        for child_shape in child_shapes:
            lines.extend(shape_lines(child_shape))
        return lines

    return []


def convert_pptx_to_md(source: str | Path, target: str | Path) -> Path:
    source_path = Path(source)
    target_path = Path(target)

    if not source_path.exists():
        raise FileNotFoundError(f"PPTX source file not found: {source_path}")

    presentation = Presentation(source_path)
    output_lines: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        output_lines.append(f"## Slide {slide_index}")
        output_lines.append("")

        title_shape = slide.shapes.title
        title_text = collapse_text(title_shape.text) if title_shape is not None else ""
        if title_text:
            output_lines.append(f"### {title_text}")
            output_lines.append("")

        slide_lines: list[str] = []
        for shape in slide.shapes:
            if title_shape is not None and shape == title_shape:
                continue
            slide_lines.extend(shape_lines(shape))

        if slide_lines:
            output_lines.extend(slide_lines)
            output_lines.append("")

    markdown = "\n".join(output_lines).strip()
    if markdown:
        markdown += "\n"

    target_path.parent.mkdir(parents=True, exist_ok=True)
    target_path.write_text(markdown, encoding="utf-8")
    return target_path


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Convert a PowerPoint deck to markdown.")
    parser.add_argument("source", help="Path to the source .pptx file")
    parser.add_argument("target", help="Path to the target .md file")
    args = parser.parse_args(argv)

    try:
        convert_pptx_to_md(args.source, args.target)
    except Exception as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 1

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
